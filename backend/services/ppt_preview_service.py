from __future__ import annotations

import asyncio
import io
import shutil
import tempfile
import threading
from datetime import datetime
from pathlib import Path
from typing import Any
from urllib.parse import urlparse

import httpx
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.enum.dml import MSO_FILL
from pptx.enum.shapes import MSO_SHAPE_TYPE

from core.database import SessionLocal
from models.resource import LearningResource

BACKEND_DIR = Path(__file__).resolve().parents[1]
STATIC_DIR = BACKEND_DIR / "static"
PPT_DIR = STATIC_DIR / "ppt"
PREVIEW_DIR = STATIC_DIR / "ppt_preview"

PREVIEW_MODE = "python_low_fidelity"
PREVIEW_WARNING = "当前为低保真预览，完整样式请下载 PPTX 查看"
CANVAS_WIDTH = 1280
MIN_CANVAS_HEIGHT = 720

_ACTIVE_TASKS: set[int] = set()
_TASK_LOCK = threading.Lock()
_FONT_CACHE: dict[tuple[int, bool], ImageFont.FreeTypeFont | ImageFont.ImageFont] = {}


def _utc_now() -> str:
    return datetime.utcnow().isoformat()


def _preview_payload(
    status: str,
    images: list[str] | None = None,
    total: int = 0,
    error: str | None = None,
) -> dict[str, Any]:
    return {
        "status": status,
        "images": images or [],
        "total": total,
        "mode": PREVIEW_MODE,
        "warning": PREVIEW_WARNING,
        "error": error,
        "generated_at": _utc_now() if status in {"ready", "failed"} else None,
    }


def _with_preview(content: Any, preview: dict[str, Any]) -> dict[str, Any]:
    data = dict(content or {}) if isinstance(content, dict) else {}
    data["preview"] = preview
    return data


def _query_resource(db, resource_id: int, user_id: str | None = None) -> LearningResource | None:
    query = db.query(LearningResource).filter(LearningResource.id == resource_id)
    if user_id:
        query = query.filter(LearningResource.user_id == user_id)
    return query.first()


def _set_preview(resource_id: int, preview: dict[str, Any]) -> dict[str, Any]:
    db = SessionLocal()
    try:
        resource = _query_resource(db, resource_id)
        if not resource:
            return preview
        resource.content = _with_preview(resource.content, preview)
        db.commit()
        return preview
    finally:
        db.close()


def _safe_local_ppt_path(value: str) -> Path | None:
    if not value:
        return None
    filename = Path(value).name
    if not filename.lower().endswith(".pptx"):
        return None
    path = (PPT_DIR / filename).resolve()
    ppt_root = PPT_DIR.resolve()
    if ppt_root not in path.parents and path != ppt_root:
        return None
    return path if path.exists() else None


def _ppt_path_from_static_url(url: str) -> Path | None:
    parsed = urlparse(url or "")
    path = parsed.path or url
    if not path.startswith("/static/ppt/"):
        return None
    return _safe_local_ppt_path(path)


def _prepare_source_ppt(content: dict[str, Any], work_dir: Path) -> Path:
    pptx_file = str(content.get("pptx_file") or "").strip()
    pptx_url = str(content.get("pptx_url") or "").strip()

    local_path = _safe_local_ppt_path(pptx_file) or _ppt_path_from_static_url(pptx_url)
    if local_path:
        return local_path

    if pptx_url.startswith("http://") or pptx_url.startswith("https://"):
        target = work_dir / "source.pptx"
        with httpx.Client(timeout=120.0, follow_redirects=True) as client:
            response = client.get(pptx_url)
            if response.status_code != 200:
                raise RuntimeError(f"PPTX 下载失败，HTTP {response.status_code}")
            target.write_bytes(response.content)
        return target

    raise RuntimeError("未找到可用于预览的 PPTX 文件")


def _font_candidates(bold: bool = False) -> list[str]:
    windows_fonts = Path("C:/Windows/Fonts")
    if bold:
        return [
            str(windows_fonts / "msyhbd.ttc"),
            str(windows_fonts / "simhei.ttf"),
            str(windows_fonts / "arialbd.ttf"),
            "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf",
        ]
    return [
        str(windows_fonts / "msyh.ttc"),
        str(windows_fonts / "simhei.ttf"),
        str(windows_fonts / "arial.ttf"),
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
    ]


def _load_font(size: int, bold: bool = False) -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
    key = (max(8, int(size)), bool(bold))
    cached = _FONT_CACHE.get(key)
    if cached:
        return cached
    for font_path in _font_candidates(bold):
        if Path(font_path).exists():
            try:
                font = ImageFont.truetype(font_path, key[0])
                _FONT_CACHE[key] = font
                return font
            except OSError:
                continue
    font = ImageFont.load_default()
    _FONT_CACHE[key] = font
    return font


def _rgb_tuple(value: Any) -> tuple[int, int, int] | None:
    if value is None:
        return None
    text = str(value)
    if len(text) != 6:
        return None
    try:
        return tuple(int(text[index:index + 2], 16) for index in (0, 2, 4))
    except ValueError:
        return None


def _fore_color(fill: Any) -> tuple[int, int, int] | None:
    try:
        if fill.type == MSO_FILL.SOLID:
            return _rgb_tuple(fill.fore_color.rgb)
    except Exception:
        return None
    return None


def _line_color(shape: Any) -> tuple[int, int, int] | None:
    try:
        return _rgb_tuple(shape.line.color.rgb)
    except Exception:
        return None


def _slide_background(slide: Any) -> tuple[int, int, int]:
    try:
        color = _fore_color(slide.background.fill)
        if color:
            return color
    except Exception:
        pass
    return (255, 255, 255)


def _shape_bbox(shape: Any, scale_x: float, scale_y: float) -> tuple[int, int, int, int]:
    left = int(shape.left * scale_x)
    top = int(shape.top * scale_y)
    right = int((shape.left + shape.width) * scale_x)
    bottom = int((shape.top + shape.height) * scale_y)
    return left, top, max(left + 1, right), max(top + 1, bottom)


def _text_size(draw: ImageDraw.ImageDraw, text: str, font: ImageFont.ImageFont) -> tuple[int, int]:
    box = draw.textbbox((0, 0), text, font=font)
    return box[2] - box[0], box[3] - box[1]


def _wrap_text(draw: ImageDraw.ImageDraw, text: str, font: ImageFont.ImageFont, max_width: int) -> list[str]:
    if not text:
        return []
    tokens = text.split(" ") if " " in text else list(text)
    lines: list[str] = []
    current_line = ""
    separator = " " if " " in text else ""
    for token in tokens:
        candidate_line = token if not current_line else current_line + separator + token
        candidate_width, _ = _text_size(draw, candidate_line, font)
        if candidate_width <= max_width or not current_line:
            current_line = candidate_line
        else:
            lines.append(current_line)
            current_line = token
    if current_line:
        lines.append(current_line)
    return lines


def _paragraph_font_size(paragraph: Any, fallback: int) -> int:
    for run in paragraph.runs:
        try:
            if run.font.size:
                return max(10, int(run.font.size.pt))
        except Exception:
            continue
    return fallback


def _paragraph_color(paragraph: Any) -> tuple[int, int, int]:
    for run in paragraph.runs:
        try:
            color = _rgb_tuple(run.font.color.rgb)
            if color:
                return color
        except Exception:
            continue
    return (38, 38, 38)


def _draw_text_frame(
    draw: ImageDraw.ImageDraw,
    shape: Any,
    bbox: tuple[int, int, int, int],
    scale_y: float,
) -> None:
    left, top, right, bottom = bbox
    padding = max(8, int(10 * scale_y))
    cursor_y = top + padding
    max_width = max(10, right - left - padding * 2)
    max_height = bottom - top - padding * 2

    for paragraph in shape.text_frame.paragraphs:
        text = paragraph.text.strip()
        if not text:
            cursor_y += 8
            continue
        level = int(getattr(paragraph, "level", 0) or 0)
        prefix = "• " if level > 0 else ""
        font_size = int(_paragraph_font_size(paragraph, 18) * scale_y)
        font = _load_font(font_size)
        text_color = _paragraph_color(paragraph)
        indent = level * 28
        lines = _wrap_text(draw, prefix + text, font, max_width - indent)
        line_height = max(font_size + 6, 18)
        for line in lines:
            if cursor_y + line_height > top + padding + max_height:
                return
            draw.text((left + padding + indent, cursor_y), line, font=font, fill=text_color)
            cursor_y += line_height
        cursor_y += max(4, int(4 * scale_y))


def _draw_shape_box(draw: ImageDraw.ImageDraw, shape: Any, bbox: tuple[int, int, int, int]) -> None:
    fill_color = _fore_color(shape.fill) if hasattr(shape, "fill") else None
    outline_color = _line_color(shape)
    if not fill_color and not outline_color:
        return
    draw.rectangle(bbox, fill=fill_color, outline=outline_color, width=2 if outline_color else 1)


def _draw_picture(canvas: Image.Image, draw: ImageDraw.ImageDraw, shape: Any, bbox: tuple[int, int, int, int]) -> None:
    left, top, right, bottom = bbox
    target_width = max(1, right - left)
    target_height = max(1, bottom - top)
    try:
        picture = Image.open(io.BytesIO(shape.image.blob)).convert("RGBA")
        picture.thumbnail((target_width, target_height), Image.Resampling.LANCZOS)
        paste_left = left + (target_width - picture.width) // 2
        paste_top = top + (target_height - picture.height) // 2
        canvas.paste(picture, (paste_left, paste_top), picture)
    except Exception:
        _draw_placeholder(draw, bbox, "图片预览不可用")


def _draw_placeholder(draw: ImageDraw.ImageDraw, bbox: tuple[int, int, int, int], label: str) -> None:
    left, top, right, bottom = bbox
    draw.rectangle(bbox, outline=(180, 180, 180), width=2)
    font = _load_font(18)
    text_width, text_height = _text_size(draw, label, font)
    draw.text(
        (left + max(4, (right - left - text_width) // 2), top + max(4, (bottom - top - text_height) // 2)),
        label,
        font=font,
        fill=(120, 120, 120),
    )


def _draw_table(draw: ImageDraw.ImageDraw, shape: Any, bbox: tuple[int, int, int, int]) -> None:
    left, top, right, bottom = bbox
    table = shape.table
    row_count = len(table.rows)
    col_count = len(table.columns)
    if row_count == 0 or col_count == 0:
        _draw_placeholder(draw, bbox, "表格预览不可用")
        return
    cell_width = max(1, (right - left) // col_count)
    cell_height = max(1, (bottom - top) // row_count)
    font = _load_font(14)
    for row_index, row in enumerate(table.rows):
        for col_index, cell in enumerate(row.cells):
            cell_left = left + col_index * cell_width
            cell_top = top + row_index * cell_height
            cell_box = (cell_left, cell_top, cell_left + cell_width, cell_top + cell_height)
            draw.rectangle(cell_box, outline=(170, 170, 170), width=1)
            text = cell.text.strip()
            if text:
                lines = _wrap_text(draw, text, font, max(8, cell_width - 8))[:2]
                cursor_y = cell_top + 4
                for line in lines:
                    draw.text((cell_left + 4, cursor_y), line, font=font, fill=(45, 45, 45))
                    cursor_y += 18


def _draw_low_fidelity_badge(draw: ImageDraw.ImageDraw) -> None:
    font = _load_font(18)
    label = "低保真预览"
    text_width, text_height = _text_size(draw, label, font)
    box = (16, 14, 16 + text_width + 22, 14 + text_height + 14)
    draw.rounded_rectangle(box, radius=8, fill=(255, 246, 214), outline=(232, 170, 60), width=1)
    draw.text((27, 20), label, font=font, fill=(138, 92, 0))


def _render_shape(canvas: Image.Image, draw: ImageDraw.ImageDraw, shape: Any, scale_x: float, scale_y: float) -> None:
    bbox = _shape_bbox(shape, scale_x, scale_y)
    shape_type = getattr(shape, "shape_type", None)

    if shape_type == MSO_SHAPE_TYPE.PICTURE:
        _draw_picture(canvas, draw, shape, bbox)
        return

    if getattr(shape, "has_table", False):
        _draw_table(draw, shape, bbox)
        return

    if getattr(shape, "has_chart", False):
        _draw_placeholder(draw, bbox, "图表预览不可用")
        return

    if shape_type == MSO_SHAPE_TYPE.GROUP:
        _draw_placeholder(draw, bbox, "组合元素已降级")
        return

    _draw_shape_box(draw, shape, bbox)
    if getattr(shape, "has_text_frame", False):
        try:
            _draw_text_frame(draw, shape, bbox, scale_y)
        except Exception:
            _draw_placeholder(draw, bbox, "文本预览不可用")


def _render_pptx_to_images(ppt_path: Path, resource_id: int, render_dir: Path) -> list[str]:
    presentation = Presentation(str(ppt_path))
    slide_width = int(presentation.slide_width)
    slide_height = int(presentation.slide_height)
    if slide_width <= 0 or slide_height <= 0:
        raise RuntimeError("PPTX 页面尺寸异常")

    canvas_width = CANVAS_WIDTH
    canvas_height = max(MIN_CANVAS_HEIGHT, int(canvas_width * slide_height / slide_width))
    scale_x = canvas_width / slide_width
    scale_y = canvas_height / slide_height

    render_dir.mkdir(parents=True, exist_ok=True)
    images: list[str] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        canvas = Image.new("RGB", (canvas_width, canvas_height), _slide_background(slide))
        draw = ImageDraw.Draw(canvas)
        for shape in slide.shapes:
            try:
                _render_shape(canvas, draw, shape, scale_x, scale_y)
            except Exception:
                bbox = _shape_bbox(shape, scale_x, scale_y)
                _draw_placeholder(draw, bbox, "元素预览不可用")
        _draw_low_fidelity_badge(draw)
        image_name = f"slide_{slide_index:03d}.png"
        image_path = render_dir / image_name
        canvas.save(image_path, format="PNG", optimize=True)
        images.append(f"/static/ppt_preview/{resource_id}/{image_name}")

    if not images:
        raise RuntimeError("PPTX 不包含可预览页面")
    return images


def get_ppt_preview_status(resource_id: int, user_id: str | None = None) -> dict[str, Any]:
    db = SessionLocal()
    try:
        resource = _query_resource(db, resource_id, user_id)
        if not resource:
            return {"found": False}
        if resource.resource_type != "ppt":
            return {"found": True, "ok": False, "error": "资源不是 PPT 类型", "preview": None}
        content = resource.content or {}
        preview = content.get("preview") if isinstance(content, dict) else None
        return {
            "found": True,
            "ok": True,
            "resource_id": resource.id,
            "preview": preview or _preview_payload("idle"),
        }
    finally:
        db.close()


def generate_ppt_preview(resource_id: int, force: bool = False) -> dict[str, Any]:
    db = SessionLocal()
    try:
        resource = _query_resource(db, resource_id)
        if not resource:
            raise RuntimeError("资源不存在")
        if resource.resource_type != "ppt":
            raise RuntimeError("资源不是 PPT 类型")
        content = dict(resource.content or {})
        preview = content.get("preview") or {}
        if preview.get("status") == "ready" and not force:
            return preview

        pending = _preview_payload("pending")
        resource.content = _with_preview(content, pending)
        db.commit()
    finally:
        db.close()

    PREVIEW_DIR.mkdir(parents=True, exist_ok=True)
    try:
        with tempfile.TemporaryDirectory(prefix=f"ppt_preview_{resource_id}_", dir=str(PREVIEW_DIR)) as temp_dir:
            work_dir = Path(temp_dir)
            source_path = _prepare_source_ppt(content, work_dir)
            staged_dir = work_dir / "slides"
            images = _render_pptx_to_images(source_path, resource_id, staged_dir)

            final_dir = PREVIEW_DIR / str(resource_id)
            if final_dir.exists():
                shutil.rmtree(final_dir)
            shutil.move(str(staged_dir), str(final_dir))

        ready = _preview_payload("ready", images=images, total=len(images))
        return _set_preview(resource_id, ready)
    except Exception as exc:
        failed = _preview_payload("failed", error=str(exc)[:500])
        return _set_preview(resource_id, failed)


def schedule_ppt_preview(resource_id: int, force: bool = False) -> bool:
    with _TASK_LOCK:
        if resource_id in _ACTIVE_TASKS:
            return False
        _ACTIVE_TASKS.add(resource_id)

    def runner() -> None:
        try:
            generate_ppt_preview(resource_id, force=force)
        finally:
            with _TASK_LOCK:
                _ACTIVE_TASKS.discard(resource_id)

    try:
        loop = asyncio.get_running_loop()
        loop.create_task(asyncio.to_thread(runner))
    except RuntimeError:
        threading.Thread(target=runner, daemon=True).start()
    return True


def start_ppt_preview(resource_id: int, user_id: str | None = None, force: bool = True) -> dict[str, Any]:
    db = SessionLocal()
    try:
        resource = _query_resource(db, resource_id, user_id)
        if not resource:
            return {"found": False}
        if resource.resource_type != "ppt":
            return {"found": True, "ok": False, "error": "资源不是 PPT 类型", "preview": None}

        content = dict(resource.content or {})
        preview = content.get("preview") or {}
        if preview.get("status") == "ready" and not force:
            return {"found": True, "ok": True, "resource_id": resource.id, "preview": preview}

        pending = _preview_payload("pending")
        resource.content = _with_preview(content, pending)
        db.commit()
    finally:
        db.close()

    schedule_ppt_preview(resource_id, force=force)
    return {
        "found": True,
        "ok": True,
        "resource_id": resource_id,
        "preview": pending,
    }


def cleanup_ppt_preview(resource_id: int) -> None:
    try:
        target = PREVIEW_DIR / str(int(resource_id))
    except (TypeError, ValueError):
        return
    preview_root = PREVIEW_DIR.resolve()
    resolved = target.resolve()
    if preview_root not in resolved.parents:
        return
    if resolved.exists():
        shutil.rmtree(resolved, ignore_errors=True)
