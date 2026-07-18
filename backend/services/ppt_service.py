"""LLM 输出的 PPT JSON → HTML 渲染 → Playwright 截图 → 贴图入 .pptx"""
import os
import uuid
import html
import asyncio
from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Inches

PPT_DIR = os.path.join(os.path.dirname(os.path.dirname(__file__)), "static", "ppt")
VIEWPORT = {"width": 1920, "height": 1080}
SLIDE_W = 13.333
SLIDE_H = 7.5

CSS = """
* { margin: 0; padding: 0; box-sizing: border-box; }
body {
  width: 1920px; height: 1080px; overflow: hidden;
  font-family: "Microsoft YaHei", "PingFang SC", "Noto Sans SC", "Helvetica Neue", sans-serif;
  -webkit-font-smoothing: antialiased;
}

/* ===== 封面 ===== */
.cover {
  width: 100%; height: 100%; position: relative;
  background: linear-gradient(150deg, #0d1b3e 0%, #13254a 25%, #18305e 50%, #1f3f70 75%, #254d82 100%);
  display: flex; flex-direction: column; align-items: flex-start; justify-content: center;
  padding: 0 140px; overflow: hidden;
}
.cover-bg1 { position: absolute; top: -200px; right: -100px; width: 800px; height: 800px;
  border-radius: 50%; border: 1px solid rgba(255,255,255,0.06);
  background: radial-gradient(circle, rgba(100,160,255,0.10) 0%, transparent 60%); }
.cover-bg2 { position: absolute; bottom: -150px; left: -80px; width: 600px; height: 600px;
  border-radius: 50%; border: 1px solid rgba(255,255,255,0.04);
  background: radial-gradient(circle, rgba(100,200,180,0.08) 0%, transparent 60%); }
.cover-bg3 { position: absolute; top: 50%; right: 10%; width: 300px; height: 300px;
  border-radius: 50%; border: 2px dashed rgba(255,255,255,0.06); transform: translateY(-50%); }
.cover-tag {
  font-size: 18px; font-weight: 600; color: rgba(100,180,255,0.9);
  letter-spacing: 8px; text-transform: uppercase; margin-bottom: 36px;
  position: relative; z-index: 1;
}
.cover-title {
  font-size: 78px; font-weight: 900; color: #ffffff;
  line-height: 1.25; max-width: 1400px;
  position: relative; z-index: 1;
  text-shadow: 0 6px 40px rgba(0,0,0,0.35);
}
.cover-divider {
  width: 100px; height: 5px; border-radius: 3px;
  background: linear-gradient(90deg, #4fc3f7, #81c784);
  margin: 40px 0; position: relative; z-index: 1;
}
.cover-sub {
  font-size: 26px; color: rgba(255,255,255,0.55);
  letter-spacing: 4px; position: relative; z-index: 1;
}

/* ===== 内容页 ===== */
.slide {
  width: 100%; height: 100%; background: #fafbfd;
  display: flex; position: relative;
  background-image: radial-gradient(circle at 95% 5%, rgba(37,99,163,0.04) 0%, transparent 50%),
                    radial-gradient(circle at 5% 95%, rgba(37,99,163,0.03) 0%, transparent 50%);
}
.slide-accent {
  width: 14px; height: 100%; flex-shrink: 0;
  background: linear-gradient(180deg, #1e3a5f 0%, #2565a0 40%, #4fc3f7 100%);
}
.slide-body { flex: 1; display: flex; flex-direction: column; padding: 70px 90px 60px 80px; }
.slide-title-area {
  display: flex; align-items: center; gap: 24px; margin-bottom: 56px;
}
.slide-title-icon {
  width: 56px; height: 56px; border-radius: 16px;
  background: linear-gradient(135deg, #2565a0, #4fc3f7);
  display: flex; align-items: center; justify-content: center;
  color: #fff; font-size: 28px; font-weight: 800; flex-shrink: 0;
}
.slide-title {
  font-size: 50px; font-weight: 800; color: #0d1b3e;
  line-height: 1.2; letter-spacing: -0.5px;
}
.slide-points { list-style: none; flex: 1; display: flex; flex-direction: column; gap: 4px; }
.slide-points li { display: flex; align-items: flex-start; gap: 24px; }
.point-num {
  width: 52px; height: 52px; border-radius: 14px; flex-shrink: 0;
  display: flex; align-items: center; justify-content: center;
  color: #fff; font-size: 24px; font-weight: 700; margin-top: 2px;
}
.point-num.c0 { background: linear-gradient(135deg, #2565a0, #4fc3f7); }
.point-num.c1 { background: linear-gradient(135deg, #00796b, #26a69a); }
.point-num.c2 { background: linear-gradient(135deg, #c62828, #ef5350); }
.point-num.c3 { background: linear-gradient(135deg, #e65100, #ff7043); }
.point-num.c4 { background: linear-gradient(135deg, #6a1b9a, #ab47bc); }
.point-num.c5 { background: linear-gradient(135deg, #1565c0, #42a5f5); }
.point-card {
  flex: 1; background: #ffffff; border-radius: 12px;
  padding: 10px 24px; border: 1px solid rgba(0,0,0,0.04);
  box-shadow: 0 2px 8px rgba(0,0,0,0.03);
}
.point-text { font-size: 27px; color: #1a1a2e; line-height: 52px; }
.slide-notes {
  width: 100%; padding: 18px 32px; margin-top: 20px;
  background: linear-gradient(90deg, rgba(37,99,163,0.06), rgba(37,99,163,0.02));
  border-left: 5px solid #2565a0; border-radius: 4px;
  font-size: 20px; color: #4a5568;
}
.notes-label { font-weight: 700; color: #2565a0; margin-right: 12px; }
.page-num {
  position: absolute; bottom: 28px; right: 64px;
  font-size: 17px; color: #b0b8c4; font-weight: 500;
}
"""

COVER_HTML = """<!DOCTYPE html><html><head><meta charset="utf-8"><style>{css}</style></head>
<body><div class="cover">
<div class="cover-bg1"></div><div class="cover-bg2"></div><div class="cover-bg3"></div>
<div class="cover-tag">ZHITU · AI COURSEWARE</div>
<div class="cover-title">{title}</div>
<div class="cover-divider"></div>
<div class="cover-sub">智途 · AI 个性化学习课件</div>
</div></body></html>"""

SLIDE_HTML = """<!DOCTYPE html><html><head><meta charset="utf-8"><style>{css}</style></head>
<body><div class="slide">
<div class="slide-accent"></div>
<div class="slide-body">
<div class="slide-title-area">
<div class="slide-title-icon">{icon_num}</div>
<div class="slide-title">{title}</div>
</div>
<ul class="slide-points">{items}</ul>
{notes_html}
</div>
<span class="page-num">{page_num}</span>
</div></body></html>"""

ITEM_HTML = '<li><span class="point-num c{color_i}">{i}</span><div class="point-card"><span class="point-text">{text}</span></div></li>'
NOTES_HTML = '<div class="slide-notes"><span class="notes-label">备注</span>{notes}</div>'


def _build_pptx(ppt_data: dict) -> dict:
    os.makedirs(PPT_DIR, exist_ok=True)

    slide_htmls = []

    with sync_playwright() as p:
        browser = p.chromium.launch()
        page = browser.new_page(viewport=VIEWPORT)

        cover_html = COVER_HTML.format(
            css=CSS,
            title=html.escape(ppt_data.get("title", "课件")),
        )
        slide_htmls.append(cover_html)
        page.set_content(cover_html, wait_until="networkidle")
        cover_img = page.screenshot(full_page=False, type="png")

        slide_imgs = []
        slides = ppt_data.get("slides", [])
        for idx, s in enumerate(slides):
            items_html = ""
            for i, point in enumerate(s.get("content", []), 1):
                color_i = (i - 1) % 6
                items_html += ITEM_HTML.format(i=i, color_i=color_i, text=html.escape(point))
            notes_html = ""
            if s.get("notes"):
                notes_html = NOTES_HTML.format(notes=html.escape(s["notes"]))
            slide_html = SLIDE_HTML.format(
                css=CSS,
                icon_num=idx + 1,
                title=html.escape(s.get("title", "")),
                items=items_html,
                notes_html=notes_html,
                page_num=f"{idx + 2} / {len(slides) + 1}",
            )
            slide_htmls.append(slide_html)
            page.set_content(slide_html, wait_until="networkidle")
            slide_imgs.append(page.screenshot(full_page=False, type="png"))

        page.close()
        browser.close()

    prefix = uuid.uuid4().hex[:8]
    img_dir = os.path.join(PPT_DIR, "img")
    os.makedirs(img_dir, exist_ok=True)

    all_imgs = [cover_img] + slide_imgs
    img_paths = []
    for index, img in enumerate(all_imgs):
        path = os.path.join(img_dir, f"{prefix}_{index}.png")
        with open(path, "wb") as f:
            f.write(img)
        img_paths.append(path)

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    blank_layout = prs.slide_layouts[6]
    for img_path in img_paths:
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(img_path, Inches(0), Inches(0),
                                 Inches(SLIDE_W), Inches(SLIDE_H))

    pptx_path = os.path.join(PPT_DIR, f"{prefix}.pptx")
    prs.save(pptx_path)

    for p in img_paths:
        try:
            os.remove(p)
        except OSError:
            pass

    return {"pptx_path": pptx_path, "slide_htmls": slide_htmls}


async def generate_pptx(ppt_data: dict) -> dict:
    return await asyncio.to_thread(_build_pptx, ppt_data)
