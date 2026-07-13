"""
Skill 系统：为 Agent 面板提供可插拔的能力模块。

每个 Skill 是一个独立的能力单元，由 plan 节点根据任务自动选择调用。
Skill 执行结果通过 SSE 事件流实时推送至前端。
"""

import asyncio
import copy
import json
import re
import uuid
from abc import ABC, abstractmethod
from dataclasses import dataclass, field
from typing import Any
from services.resource_title_service import build_resource_title


@dataclass
class SkillResult:
    """Skill 执行结果"""
    success: bool = True
    data: dict = field(default_factory=dict)
    summary: str = ""
    error: str = ""


class BaseSkill(ABC):
    """Skill 基类，所有 skill 需继承此类"""
    name: str = ""
    description: str = ""  # 供 LLM plan 节点选择参考
    icon: str = "🔧"

    @abstractmethod
    async def execute(self, context: dict, workflow_outputs: list) -> SkillResult:
        """
        执行 skill。
        
        Args:
            context: 包含 user_message, user_id, all_modules_data 等上下文
            workflow_outputs: 当前 workflow 输出列表，skill 可向其中追加步骤事件
        
        Returns:
            SkillResult 包含执行结果
        """
        ...

    _SKILL_AGENT_NAMES: dict = {
        "deep_search": "搜索智能体", "code_analysis": "代码智能体",
        "mindmap_gen": "导图智能体", "quiz_gen": "出题智能体", "video_search": "视频智能体",
        "ppt_gen": "课件智能体", "article_gen": "文章智能体",
        "code_gen": "代码/动画智能体", "practice_case": "实践案例智能体",
    }

    def emit_step(self, workflow_outputs: list, status: str, title: str, data: dict, step_id: str | None = None) -> str:
        """向 workflow_outputs 追加一个 skill step 事件，并实时推送到 SSE 队列"""
        sid = step_id or str(uuid.uuid4())[:8]
        q: asyncio.Queue | None = getattr(self, "_sse_queue", None)
        event = {
            "type": "step",
            "step_type": "skill",
            "step_id": sid,
            "status": status,
            "title": title,
            "agent_name": self._SKILL_AGENT_NAMES.get(self.name, self.name),
            "_live_pushed": q is not None,
            "data": {
                "skill_name": self.name,
                "skill_icon": self.icon,
                **data,
            },
        }
        workflow_outputs.append(event)
        # 实时推送：如果 context 中有 SSE 队列，立即 put
        if q is not None:
            try:
                q.put_nowait(event)
            except asyncio.QueueFull:
                pass
        return sid

    def emit_token(self, step_id: str, delta: str):
        """向当前 skill step 实时追加 token 内容。"""
        if not delta:
            return
        q: asyncio.Queue | None = getattr(self, "_sse_queue", None)
        if q is None:
            return
        try:
            q.put_nowait({"type": "token", "step_id": step_id, "delta": delta})
            session_id = getattr(self, "_session_id", "")
            if session_id:
                from core.sse_registry import mark_live_token_step
                mark_live_token_step(session_id, step_id)
        except asyncio.QueueFull:
            pass


# ============================================================
# Skill Registry
# ============================================================

_SKILL_REGISTRY: dict[str, BaseSkill] = {}


def register_skill(skill: BaseSkill):
    """注册一个 skill 实例"""
    _SKILL_REGISTRY[skill.name] = skill


def get_skill(name: str) -> BaseSkill | None:
    """按名称获取 skill"""
    if not _SKILL_REGISTRY:
        init_skills()
    skill = _SKILL_REGISTRY.get(name)
    return copy.copy(skill) if skill else None


def get_all_skills() -> dict[str, BaseSkill]:
    """获取所有已注册 skills"""
    if not _SKILL_REGISTRY:
        init_skills()
    return _SKILL_REGISTRY.copy()


def get_skills_description() -> str:
    """生成 skills 列表描述，供 LLM plan 节点使用"""
    if not _SKILL_REGISTRY:
        init_skills()
    lines = []
    for name, skill in _SKILL_REGISTRY.items():
        lines.append(f"- {name}: {skill.description}")
    return "\n".join(lines)


def make_draft_resource(
    resource_type: str,
    title: str,
    content,
    course_name: str | None = None,
    knowledge_points: list[str] | None = None,
    kp_weights: dict | None = None,
    course_bindings: list[dict] | None = None,
) -> dict:
    final_title = title
    if not final_title or final_title in {f"{resource_type}_resource", "code_resource"}:
        final_title = build_resource_title(
            resource_type,
            content,
            fallback_text=title,
            course_name=course_name,
            knowledge_points=knowledge_points or [],
        )
    return {
        "client_draft_id": uuid.uuid4().hex,
        "resource_type": resource_type,
        "title": final_title,
        "content": content if isinstance(content, dict) else {"text": content},
        "course_name": course_name,
        "knowledge_points": knowledge_points or [],
        "kp_weights": kp_weights or {},
        "course_bindings": course_bindings or [],
        "save_required": True,
    }


_TOPIC_STOP_PHRASES = (
    "请帮我", "帮我", "请给我", "给我", "请", "生成一个", "生成", "制作一个", "制作",
    "创建一个", "创建", "设计一个", "设计", "实现一个", "实现", "写一个", "写个", "开发一个",
    "开发", "可视化动画", "可视化演示", "交互式动画", "交互动画", "动画演示", "动态演示",
    "可视化", "动画", "演示", "页面", "代码案例", "代码", "助我理解", "帮助我理解", "用于讲解",
    "讲解", "关于", "这个", "那个",
)


def _extract_topic_terms(task: str) -> list[str]:
    text = (task or "").lower()
    for phrase in _TOPIC_STOP_PHRASES:
        text = text.replace(phrase.lower(), " ")
    parts = re.split(r"[\s，。；：、,.!?！？:;（）()\[\]{}<>《》\"'“”‘’/\\|]+|(?:以及|或者|并且|与|和|及)", text)
    generic = {"主题", "内容", "知识点", "课程", "学习", "理解", "一个", "一种", "相关", "完整", "自包含"}
    terms = []
    for part in parts:
        term = part.strip("-_+")
        if len(term) < 2 or term in generic or term.isdigit():
            continue
        if term not in terms:
            terms.append(term)
    return sorted(terms, key=len, reverse=True)[:8]


def _validate_topic_alignment(task: str, content: str) -> tuple[bool, str]:
    terms = _extract_topic_terms(task)
    if not terms:
        return True, ""
    output = (content or "").lower()
    if any(term in output for term in terms):
        return True, ""
    return False, f"生成内容未体现规划主题关键词：{', '.join(terms[:4])}"


def _clean_html_output(content: str) -> str:
    cleaned = re.sub(r'^```html?\s*\n?', '', (content or "").strip(), flags=re.IGNORECASE)
    cleaned = re.sub(r'\n?```\s*$', '', cleaned)
    start = re.search(r'<(!DOCTYPE\s+html|html[\s>])', cleaned, re.IGNORECASE)
    end = re.search(r'</html\s*>', cleaned, re.IGNORECASE)
    if start and end:
        return cleaned[start.start():end.end()].strip()
    if end:
        return cleaned[:end.end()].strip()
    first_lt = cleaned.find('<')
    last_gt = cleaned.rfind('>')
    if first_lt >= 0 and last_gt > first_lt:
        return cleaned[first_lt:last_gt + 1].strip()
    return cleaned


def infer_draft_resource_binding(
    text: str,
    course_name: str | None = None,
    knowledge_points: list[str] | None = None,
) -> dict:
    from services.kp_service import infer_resource_tags

    graph_tags = infer_resource_tags(
        text or "",
        course_name=course_name,
        knowledge_points=knowledge_points or [],
    )
    return {
        "course_name": graph_tags.get("course_name"),
        "knowledge_points": graph_tags.get("knowledge_points") or [],
        "kp_weights": graph_tags.get("kp_weights") or {},
        "course_bindings": graph_tags.get("course_bindings") or [],
    }


def _has_mistake_task_intent(message: str) -> bool:
    return any(keyword in (message or "") for keyword in ["错题", "错因", "错误题", "薄弱题", "错题本"])


def _mistake_knowledge_points(mistakes: dict | None) -> list[str]:
    points: list[str] = []
    for item in (mistakes or {}).get("recent") or []:
        for point in item.get("knowledge_points") or []:
            if point and point not in points:
                points.append(str(point))
    return points[:10]


# ============================================================
# 具体 Skill 实现
# ============================================================


class DeepSearchSkill(BaseSkill):
    """深度搜索 - 多轮搜索 + 网页抓取，获取更详尽的信息"""
    name = "deep_search"
    description = "深度搜索：对任务进行多轮互联网搜索并抓取关键网页内容，适用于需要全面资料的任务"
    icon = "🔍"

    async def execute(self, context: dict, workflow_outputs: list) -> SkillResult:
        from agents.tools import tavily_search, web_scrape

        user_message = context.get("user_message", "")
        ad = context.get("all_modules_data", {})
        search_keywords = ad.get("search_keywords", user_message)

        step_id = self.emit_step(workflow_outputs, "running", f"深度搜索: {search_keywords[:30]}", {
            "sub_steps": ["⏳ 第一轮关键词搜索中..."],
        })

        # 第一轮：关键词搜索
        result1 = await tavily_search(search_keywords, max_results=5)
        sub_steps = [f"✅ 关键词搜索完成，获取 {len(result1.get('results', []))} 条结果", "⏳ 第二轮补充搜索中..."]
        self.emit_step(workflow_outputs, "running", f"深度搜索: {search_keywords[:30]}", {"sub_steps": sub_steps}, step_id)

        # 第二轮：原始问题搜索
        result2 = await tavily_search(user_message, max_results=3)
        sub_steps[-1] = f"✅ 补充搜索完成，获取 {len(result2.get('results', []))} 条结果"

        # 合并去重
        seen_urls = set()
        all_results = []
        for r in result1.get("results", []) + result2.get("results", []):
            url = r.get("url", "")
            if url and url not in seen_urls:
                seen_urls.add(url)
                all_results.append(r)

        # 抓取前2个网页获取详细内容
        scraped_contents = []
        for r in all_results[:2]:
            url = r.get("url", "")
            if url:
                sub_steps.append(f"⏳ 抓取网页: {r.get('title', url)[:40]}...")
                self.emit_step(workflow_outputs, "running", f"深度搜索: {search_keywords[:30]}", {"sub_steps": sub_steps}, step_id)
                scrape_result = await web_scrape(url)
                content = scrape_result.get("content", "")
                if content and len(content) > 50:
                    scraped_contents.append({"url": url, "title": r.get("title", ""), "content": content[:1000]})
                    sub_steps[-1] = f"✅ 抓取完成: {r.get('title', url)[:40]}"

        answer = result1.get("answer", "") or result2.get("answer", "")

        self.emit_step(workflow_outputs, "completed", f"深度搜索: {search_keywords[:30]}", {
            "content": f"搜索完成，共获取 {len(all_results)} 条结果，抓取 {len(scraped_contents)} 个网页",
            "sub_steps": sub_steps,
        }, step_id)

        return SkillResult(
            success=True,
            data={
                "search_results": all_results,
                "scraped_contents": scraped_contents,
                "answer": answer,
            },
            summary=f"深度搜索完成: {len(all_results)} 条结果, {len(scraped_contents)} 个网页内容",
        )


class CodeAnalysisSkill(BaseSkill):
    """代码分析与生成"""
    name = "code_analysis"
    description = "代码分析与生成：根据任务需求生成、解释或优化代码，适用于编程相关任务"
    icon = "💻"

    async def execute(self, context: dict, workflow_outputs: list) -> SkillResult:
        from core.llm_client import chat_completion

        user_message = context.get("user_message", "")
        ad = context.get("all_modules_data", {})
        code_lang = ad.get("code_lang", "python")
        code_desc = ad.get("code_desc", user_message)
        search_answer = ad.get("search_result", {}).get("answer", "") if ad.get("search_result") else ""

        step_id = self.emit_step(workflow_outputs, "running", f"代码生成 ({code_lang})", {
            "sub_steps": ["⏳ 正在调用模型生成代码..."],
        })

        system_prompt = "你是一个代码专家。根据任务需求生成高质量、可运行的代码。包含详细注释说明逻辑。"
        user_prompt = f"""任务：{code_desc}
上下文：{user_message}
参考信息：{search_answer[:500]}
语言：{code_lang}

请生成完整的可运行代码，包含注释。只输出代码，不要markdown标记。"""

        try:
            resp = await chat_completion([
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt},
            ], temperature=0.3)
            code_content = resp.choices[0].message.content.strip()
            # 清理 markdown 代码块标记
            if code_content.startswith("```"):
                lines = code_content.split("\n")
                code_content = "\n".join(lines[1:-1] if lines[-1].strip() == "```" else lines[1:])

            self.emit_step(workflow_outputs, "completed", f"代码生成 ({code_lang})", {
                "content": code_content,
                "sub_steps": [f"✅ 已生成 {code_lang} 代码"],
                "language": code_lang,
            }, step_id)

            return SkillResult(
                success=True,
                data={"code": code_content, "language": code_lang},
                summary=f"代码生成完成 ({code_lang})",
            )
        except Exception as e:
            self.emit_step(workflow_outputs, "completed", f"代码生成 ({code_lang})", {
                "content": f"代码生成失败: {str(e)}",
                "sub_steps": [f"❌ 错误: {str(e)}"],
            }, step_id)
            return SkillResult(success=False, error=str(e))


class MindmapSkill(BaseSkill):
    """思维导图生成"""
    name = "mindmap_gen"
    description = "思维导图生成：将知识点整理为结构化思维导图，适用于需要知识梳理和体系化的任务"
    icon = "🧠"

    async def execute(self, context: dict, workflow_outputs: list) -> SkillResult:
        from core.llm_client import chat_completion

        user_message = context.get("user_message", "")

        step_id = self.emit_step(workflow_outputs, "running", "生成思维导图", {
            "sub_steps": ["⏳ 正在调用模型生成知识结构..."],
        })

        prompt = f"""你是大学学科知识体系专家。请根据主题生成一份用于径向树图(Radial Tree)可视化的JSON知识结构。

主题：{user_message}

【核心原则】
- 每个节点必须有教育深度：name 是短标签(4-12字)，detail 是解释/定义/公式(10-40字)
- 树图只显示 name，detail 用于鼠标悬停弹窗展示
- 禁止只写空洞词典条目（如 "定义：..."），要包含原理、关系、应用等实际知识

【结构要求】
- 根节点 name 为主题精简名(5-10字)，不设 detail
- 一级分支 4-7 个，覆盖不同的知识维度（如概念、原理、方法、应用、工具、关联学科等）
- 每个分支下 2-4 个子节点，最多 3 层
- 叶节点 detail 要特别充实

JSON格式：
{{
  "name": "根节点短名",
  "children": [
    {{
      "name": "分支短标签",
      "detail": "该分支的核心知识说明",
      "children": [
        {{
          "name": "知识点短标签",
          "detail": "具体知识：包含定义、公式、关键参数或一句话示例"
        }}
      ]
    }}
  ]
}}

只返回JSON，不要其他内容。"""

        try:
            resp = await chat_completion([
                {"role": "user", "content": prompt},
            ], temperature=0.55)
            raw = resp.choices[0].message.content.strip()
            if raw.startswith("```"):
                raw = raw.strip("`").strip()
                if raw.startswith("json"):
                    raw = raw[4:].strip()

            tree_data = json.loads(raw)

            self.emit_step(workflow_outputs, "completed", "生成思维导图", {
                "content": json.dumps(tree_data, ensure_ascii=False),
                "sub_steps": ["✅ 径向树图数据生成完成"],
                "render_type": "radial_tree",
            }, step_id)

            return SkillResult(
                success=True,
                data={"tree": tree_data, "type": "mindmap"},
                summary="思维导图生成完成",
            )
        except Exception as e:
            self.emit_step(workflow_outputs, "completed", "生成思维导图", {
                "content": f"生成失败: {str(e)}",
                "sub_steps": [f"❌ 错误: {str(e)}"],
            }, step_id)
            return SkillResult(success=False, error=str(e))


class ArticleGenSkill(BaseSkill):
    """学习文章生成 — 委托给 ContentGenAgent"""
    name = "article_gen"
    description = "文章生成：根据知识点生成个性化学习文章，适用于需要详细讲解某个概念的任务"
    icon = "📄"

    async def execute(self, context: dict, workflow_outputs: list) -> SkillResult:
        from agents.content_gen_agent import ContentGenAgent
        from agents.base import AgentState

        user_message = context.get("user_message", "")
        user_id = context.get("user_id", "")
        persist = context.get("persist", True) is not False

        step_id = self.emit_step(workflow_outputs, "running", "生成学习文章", {
            "sub_steps": ["⏳ 正在调用模型生成文章..."],
        })

        try:
            state = AgentState(
                user_id=user_id,
                user_message=user_message,
                resource_type="article",
                profile=context.get("profile"),
                profile_context=context.get("profile_text"),
                persist=False,
            )
            agent = ContentGenAgent()
            await agent._generate_article(state)

            resp = json.loads(state.get("response", "{}"))
            content = resp.get("content", "")

            self.emit_step(workflow_outputs, "completed", "生成学习文章", {
                "content": content,
                "sub_steps": ["✅ 文章生成完成"],
                "draft_resource": state.get("draft_resource"),
            }, step_id)

            return SkillResult(success=True, data={"article": content, "type": "article", "draft_resource": state.get("draft_resource")}, summary="文章生成完成")
        except Exception as e:
            self.emit_step(workflow_outputs, "completed", "生成学习文章", {
                "content": f"生成失败: {str(e)}",
                "sub_steps": [f"❌ {str(e)}"],
            }, step_id)
            return SkillResult(success=False, error=str(e))


class CodeGenSkill(BaseSkill):
    """代码案例生成 — 委托给 ContentGenAgent"""
    name = "code_gen"
    description = "代码案例生成：生成带注释的可运行代码示例或算法可视化动画（HTML+JS自包含页面），适用于编程概念教学、算法步骤演示、数据结构可视化等任务"
    icon = "💡"

    async def execute(self, context: dict, workflow_outputs: list) -> SkillResult:
        from agents.content_gen_agent import CODE_PROMPT, ContentGenAgent
        from agents.base import AgentState
        from core.llm_client import chat_completion
        from services.safety_service import check_text, hallu_rules

        user_message = context.get("user_message", "")
        user_id = context.get("user_id", "")
        ad = context.get("all_modules_data", {})
        code_lang = ad.get("code_lang", "python")
        code_desc = (ad.get("code_desc") or user_message or "").strip()
        generation_task = code_desc or user_message
        binding_text = "\n".join([
            str(user_message or ""),
            str(generation_task or ""),
            str(ad.get("search_keywords") or ""),
        ])
        graph_binding = infer_draft_resource_binding(
            binding_text,
            course_name=context.get("course_name") or ad.get("course_name"),
            knowledge_points=context.get("knowledge_points") or ad.get("knowledge_points") or [],
        )

        # 检测可视化意图
        VIZ_KEYWORDS = ["可视化", "动画", "演示", "visuali", "animation", "animate", "步骤展示", "动态展示"]
        intent_text = f"{user_message}\n{generation_task}".lower()
        is_viz = any(kw in intent_text for kw in VIZ_KEYWORDS)

        sub_steps = [f"⏳ 正在识别{'可视化动画' if is_viz else '代码案例'}生成需求..."]
        step_id = self.emit_step(workflow_outputs, "running", "生成代码案例", {
            "content": "",
            "sub_steps": sub_steps,
            "progress": 10,
            "current_phase": "需求识别",
            "progress_note": "正在分析任务描述、学生画像和输出形式",
            "progress_indeterminate": False,
            "progress_label": f"1/{5 if is_viz else 4} 阶段",
            "language": "html" if is_viz else code_lang,
            "streaming_code": True,
        })

        try:
            if is_viz:
                profile = context.get("profile")
                profile_text = f"专业：{getattr(profile, 'major', '')}，年级：{getattr(profile, 'grade', '')}" if profile else ""
                sub_steps[-1] = "✅ 已识别为可视化动画任务"
                sub_steps.append("⏳ 正在调用模型生成自包含 HTML/CSS/JS...")
                self.emit_step(workflow_outputs, "running", "生成代码案例", {
                    "sub_steps": sub_steps,
                    "progress": 35,
                    "current_phase": "动画代码生成",
                    "progress_note": "模型正在生成动画页面、交互控件和演示步骤",
                    "progress_indeterminate": True,
                    "progress_label": "2/5 阶段",
                    "language": "html",
                    "streaming_code": True,
                }, step_id)
                prompt = f"""你是教学可视化与前端开发专家。请生成一个自包含 HTML 文件（内嵌 CSS + JS），用交互动画讲解指定主题。

用户原始需求：{user_message}
必须实现的主题：{generation_task}
学生背景：{profile_text or '未知'}

输出要求：
- 只输出完整 HTML，不要 Markdown 代码块，不要解释文字。
- 必须以 <!DOCTYPE html> 或 <html 开头，并以 </html> 结尾。
- 不依赖 CDN、图片或外部文件。
- 页面标题、概念说明、变量命名、动画步骤必须全部围绕“必须实现的主题”。
- 禁止替换成其他常见示例，例如 K-Means 聚类、排序算法、汉诺塔、二叉树遍历，除非它们就是必须实现的主题。
- 页面包含：标题、概念说明、图例、动画演示区、当前步骤说明、进度指示、上一步/下一步/自动播放/重置按钮。
- JavaScript 预生成 steps 数组，每一步包含状态快照、说明文本和高亮元素；用 renderStep(index) 统一更新页面。
- 交互支持：按钮切换步骤、自动播放、重置；键盘 ←/→ 切换步骤，空格播放/暂停，R 重置。
- 样式现代、宽屏友好，主容器 max-width 不小于 900px；移动端保持可读。
- 代码尽量简洁，优先保证首屏可运行、动画稳定、交互清楚。"""

                async def generate_visual_html(current_prompt: str, emit_tokens: bool) -> str:
                    content_parts: list[str] = []
                    stream = await chat_completion(
                        [{"role": "user", "content": current_prompt}],
                        temperature=0.3,
                        stream=True,
                    )
                    async for chunk in stream:
                        delta = chunk.choices[0].delta.content if chunk.choices else ""
                        if not delta:
                            continue
                        content_parts.append(delta)
                        if emit_tokens:
                            self.emit_token(step_id, delta)
                    return _clean_html_output("".join(content_parts))

                content = await generate_visual_html(prompt, emit_tokens=True)
                if not content:
                    raise RuntimeError("模型未返回可视化动画代码")
                aligned, alignment_error = _validate_topic_alignment(generation_task, content)
                if not aligned:
                    sub_steps[-1] = f"⚠️ {alignment_error}，正在按原主题重试"
                    self.emit_step(workflow_outputs, "running", "生成代码案例", {
                        "sub_steps": sub_steps,
                        "progress": 65,
                        "current_phase": "主题纠偏",
                        "progress_note": "检测到生成内容可能偏题，正在严格按规划主题重新生成",
                        "progress_indeterminate": True,
                        "progress_label": "3/6 阶段",
                        "language": "html",
                        "streaming_code": False,
                    }, step_id)
                    retry_prompt = prompt + f"""

上一版生成内容未通过主题一致性检查：{alignment_error}
本次必须在页面标题、概念说明和动画步骤中明确出现并讲解“{generation_task}”，不得改成其他示例主题。"""
                    content = await generate_visual_html(retry_prompt, emit_tokens=False)
                    aligned, alignment_error = _validate_topic_alignment(generation_task, content)
                    if not aligned:
                        raise RuntimeError(f"可视化动画主题校验失败：{alignment_error}")

                sub_steps[-1] = "✅ 动画代码生成完成并通过主题校验"
                sub_steps.append("⏳ 正在校验 HTML 结构并准备预览...")
                self.emit_step(workflow_outputs, "running", "生成代码案例", {
                    "sub_steps": sub_steps,
                    "progress": 75,
                    "current_phase": "结构校验",
                    "progress_note": "正在移除 Markdown 包裹、截取有效 HTML、准备前端预览",
                    "progress_indeterminate": False,
                    "progress_label": "3/5 阶段",
                    "streaming_code": False,
                    "language": "html",
                }, step_id)
                if not re.search(r'<html[\s>]', content, re.IGNORECASE) or not re.search(r'</html\s*>', content, re.IGNORECASE):
                    raise RuntimeError("模型返回的可视化动画不是完整 HTML 页面")

                sub_steps[-1] = "✅ HTML 结构校验完成"
                sub_steps.append("⏳ 正在生成预览卡片...")
                self.emit_step(workflow_outputs, "running", "生成代码案例", {
                    "sub_steps": sub_steps,
                    "progress": 90,
                    "current_phase": "准备预览",
                    "progress_note": "正在把动画页面交给前端沙箱预览",
                    "progress_indeterminate": False,
                    "progress_label": "4/5 阶段",
                    "language": "html",
                    "streaming_code": False,
                }, step_id)

                draft_resource = make_draft_resource(
                    "anime",
                    f"可视化动画：{generation_task[:40]}",
                    {"code": content, "language": "html"},
                    **graph_binding,
                )
                self.emit_step(workflow_outputs, "completed", "生成代码案例", {
                    "content": content,
                    "sub_steps": sub_steps[:-1] + ["✅ 可视化动画生成完成"],
                    "progress": 100,
                    "current_phase": "生成完成",
                    "progress_note": "可视化动画已生成，可直接预览或查看代码",
                    "progress_indeterminate": False,
                    "progress_label": "5/5 阶段",
                    "language": "html",
                    "streaming_code": False,
                    "draft_resource": draft_resource,
                }, step_id)
                return SkillResult(
                    success=True,
                    data={"type": "anime", "task_desc": generation_task, "code": content, "draft_resource": draft_resource},
                    summary=f"可视化动画生成完成：{generation_task[:60]}",
                )

            # 非可视化：走原有流程
            sub_steps[-1] = "✅ 已识别为代码案例任务"
            sub_steps.append("⏳ 正在调用模型生成代码案例...")
            self.emit_step(workflow_outputs, "running", "生成代码案例", {
                "sub_steps": sub_steps,
                "progress": 35,
                "current_phase": "代码生成",
                "progress_note": "模型正在生成代码、注释和说明",
                "progress_indeterminate": True,
                "progress_label": "2/4 阶段",
                "language": code_lang,
                "streaming_code": True,
            }, step_id)
            profile = context.get("profile")
            profile_text = context.get("profile_text") or ""
            kb = getattr(profile, "knowledge_base", {}) if profile else {}
            level = "初级"
            for _name, score in (kb or {}).items():
                if score >= 0.7:
                    level = "高级"
                elif score >= 0.4 and level == "初级":
                    level = "中级"
            prompt = CODE_PROMPT.format(
                profile=profile_text or "暂无学生画像",
                topic=generation_task,
                code_lang=code_lang,
                level=level,
                hallu=hallu_rules(),
            )
            content_parts: list[str] = []
            stream = await chat_completion([{"role": "user", "content": prompt}], temperature=0.5, stream=True)
            async for chunk in stream:
                delta = chunk.choices[0].delta.content if chunk.choices else ""
                if delta:
                    content_parts.append(delta)
                    self.emit_token(step_id, delta)
            content = "".join(content_parts).strip()
            if not content:
                raise RuntimeError("模型未返回代码内容")
            sub_steps[-1] = "✅ 代码案例生成完成"
            sub_steps.append("⏳ 正在解析并准备资源草稿...")
            self.emit_step(workflow_outputs, "running", "生成代码案例", {
                "sub_steps": sub_steps,
                "progress": 85,
                "current_phase": "结果整理",
                "progress_note": "正在整理代码内容并准备可保存草稿",
                "progress_indeterminate": False,
                "progress_label": "3/4 阶段",
                "language": code_lang,
                "streaming_code": False,
            }, step_id)

            safe_content, _ = await check_text(content)
            state = AgentState(
                user_id=user_id,
                user_message=generation_task,
                resource_type="code",
                code_language=code_lang,
                profile=profile,
                profile_context=context.get("profile_text"),
                course_name=graph_binding.get("course_name"),
                knowledge_points=graph_binding.get("knowledge_points") or [],
                kp_weights=graph_binding.get("kp_weights") or {},
                persist=False,
            )
            agent = ContentGenAgent()
            agent._save_or_draft_resource(state, "code", {"code": safe_content, "language": code_lang})
            content = safe_content
            resource_type = "code"
            draft_resource = state.get("draft_resource")

            self.emit_step(workflow_outputs, "completed", "生成代码案例", {
                "content": content,
                "sub_steps": sub_steps[:-1] + ["✅ 代码案例生成完成"],
                "progress": 100,
                "current_phase": "生成完成",
                "progress_note": "代码案例已生成，可按需保存到学习资源",
                "progress_indeterminate": False,
                "progress_label": "4/4 阶段",
                "language": code_lang,
                "streaming_code": False,
                "draft_resource": draft_resource,
            }, step_id)
            return SkillResult(
                success=True,
                data={"type": resource_type, "task_desc": generation_task, "code": content, "draft_resource": draft_resource},
                summary=f"代码案例生成完成：{generation_task[:60]}",
            )

        except Exception as e:
            self.emit_step(workflow_outputs, "completed", "生成代码案例", {
                "content": f"生成失败: {str(e)}",
                "sub_steps": [f"❌ {str(e)}"],
                "progress": 100,
                "current_phase": "生成失败",
                "progress_note": str(e),
                "progress_indeterminate": False,
                "progress_label": "失败",
                "streaming_code": False,
            }, step_id)
            return SkillResult(success=False, error=str(e))


class QuizGenSkill(BaseSkill):
    """习题生成 — 委托给 ContentGenAgent"""
    name = "quiz_gen"
    description = "习题生成：根据知识点生成练习题和测验，适用于需要检验学习效果或出题的任务"
    icon = "📝"

    async def execute(self, context: dict, workflow_outputs: list) -> SkillResult:
        from agents.content_gen_agent import ContentGenAgent
        from agents.base import AgentState

        user_message = context.get("user_message", "")
        user_id = context.get("user_id", "")
        ad = context.get("all_modules_data", {})
        agent_context = ad.get("agent_context") or {}
        mistakes = agent_context.get("mistakes") or {}
        target_knowledge_points = _mistake_knowledge_points(mistakes) if _has_mistake_task_intent(user_message) else []
        if target_knowledge_points:
            from services.agent_context_service import build_mistake_prompt_context
            mistake_context = build_mistake_prompt_context(mistakes)
            user_message = (
                f"{user_message}\n\n"
                "【真实错题本上下文】\n"
                f"{mistake_context}\n\n"
                "【针对性出题要求】\n"
                f"- 必须围绕这些错题暴露出的知识点生成练习：{'、'.join(target_knowledge_points)}。\n"
                "- 题目应考察同类概念、同类推理路径或同类易错点，但不能原题照抄。\n"
                "- 每道题的 knowledge_points 必须从上述知识点中选择。"
            )

        question_count = context.get("question_count", 5)
        difficulty = context.get("difficulty", "中等")
        sub_steps = ["⏳ 正在分析题库生成要求..."]
        step_id = self.emit_step(workflow_outputs, "running", "生成练习题", {
            "sub_steps": sub_steps,
            "progress": 10,
            "current_phase": "需求分析",
            "progress_note": "正在整理知识点、题量、难度和学生画像",
            "progress_indeterminate": False,
            "progress_label": "1/5 阶段",
        })

        try:
            sub_steps[-1] = "✅ 题库要求已整理"
            sub_steps.append(f"⏳ 正在生成 {question_count} 道{difficulty}难度练习题...")
            self.emit_step(workflow_outputs, "running", "生成练习题", {
                "sub_steps": sub_steps,
                "progress": 35,
                "current_phase": "题目生成",
                "progress_note": "模型正在生成题干、选项、答案和解析",
                "progress_indeterminate": True,
                "progress_label": "2/5 阶段",
            }, step_id)
            state = AgentState(
                user_id=user_id,
                user_message=user_message,
                resource_type="quiz",
                question_count=question_count,
                difficulty=difficulty,
                code_language=ad.get("code_lang", "python"),
                knowledge_points=target_knowledge_points,
                profile=context.get("profile"),
                profile_context=context.get("profile_text"),
                persist=False,
            )
            agent = ContentGenAgent()
            await agent._generate_quiz(state)
            sub_steps[-1] = "✅ 练习题内容生成完成"
            sub_steps.append("⏳ 正在解析 JSON、校验题目结构和选项...")
            self.emit_step(workflow_outputs, "running", "生成练习题", {
                "sub_steps": sub_steps,
                "progress": 75,
                "current_phase": "结构校验",
                "progress_note": "正在确认题干、选项、答案和解析是否可展示",
                "progress_indeterminate": False,
                "progress_label": "3/5 阶段",
            }, step_id)

            resp = json.loads(state.get("response", "{}"))
            quiz_data = resp.get("content", {})
            draft_resource = state.get("draft_resource")
            total_questions = len(quiz_data.get("questions", []))
            sub_steps[-1] = f"✅ 题目结构校验完成，共 {total_questions} 题"
            sub_steps.append("⏳ 正在准备题库预览...")
            self.emit_step(workflow_outputs, "running", "生成练习题", {
                "sub_steps": sub_steps,
                "progress": 90,
                "current_phase": "准备预览",
                "progress_note": "正在把题库转换为前端可作答格式",
                "progress_indeterminate": False,
                "progress_label": "4/5 阶段",
            }, step_id)

            self.emit_step(workflow_outputs, "completed", "生成练习题", {
                "content": json.dumps(quiz_data, ensure_ascii=False, indent=2),
                "sub_steps": sub_steps[:-1] + [f"✅ 已生成 {total_questions} 道练习题"],
                "progress": 100,
                "current_phase": "生成完成",
                "progress_note": "题库已生成，可查看题目和解析",
                "progress_indeterminate": False,
                "progress_label": "5/5 阶段",
                "draft_resource": draft_resource,
            }, step_id)

            return SkillResult(
                success=True,
                data={"quiz": quiz_data, "type": "quiz", "draft_resource": draft_resource},
                summary=f"生成 {total_questions} 道练习题",
            )
        except Exception as e:
            self.emit_step(workflow_outputs, "completed", "生成练习题", {
                "content": f"生成失败: {str(e)}",
                "sub_steps": [f"❌ 错误: {str(e)}"],
                "progress": 100,
                "current_phase": "生成失败",
                "progress_note": str(e),
                "progress_indeterminate": False,
                "progress_label": "失败",
            }, step_id)
            return SkillResult(success=False, error=str(e))


class PracticeCaseSkill(BaseSkill):
    """实操案例生成 — 生成完整的实验/实践项目案例"""
    name = "practice_case"
    description = "实操案例生成：生成包含背景、目标、操作步骤、参考实现和验证要点的完整实践项目案例，适用于需要动手实践、实验练习的任务"
    icon = "🔬"

    async def execute(self, context: dict, workflow_outputs: list) -> SkillResult:
        from core.llm_client import chat_completion

        user_message = context.get("user_message", "")
        user_id = context.get("user_id", "")
        profile = context.get("profile")

        step_id = self.emit_step(workflow_outputs, "running", "生成实操案例", {
            "sub_steps": ["⏳ 正在调用模型生成实操案例..."],
        })

        profile_text = ""
        if profile:
            profile_text = f"专业：{getattr(profile, 'major', '')}，年级：{getattr(profile, 'grade', '')}，编程语言偏好：{getattr(profile, 'learning_goal', '')}"

        prompt = f"""你是一位经验丰富的实践教学专家。请为以下学习任务设计一个完整的实操案例，用Markdown格式输出。

学习任务：{user_message}
学生背景：{profile_text or '未知'}

请按以下结构输出：

## 🔬 实操案例：[案例名称]

### 📋 背景与目标
[说明本案例的应用背景和学习目标，3-5句]

### 🛠️ 环境准备
[列出所需工具、依赖、环境配置]

### 📝 实践步骤

**步骤一：[标题]**
[详细说明 + 关键代码]

**步骤二：[标题]**
[详细说明 + 关键代码]

（根据复杂度设计3-6个步骤）

### 💡 参考实现
```[语言]
[完整可运行的参考代码，含注释]
```

### ✅ 验证要点
- [ ] [验证点1]
- [ ] [验证点2]
- [ ] [验证点3]

### 🚀 拓展挑战
[1-2个进阶练习方向]

要求：步骤具体可操作，代码真实可运行，验证要点可自测。"""

        try:
            resp = await chat_completion([{"role": "user", "content": prompt}], temperature=0.4)
            content = resp.choices[0].message.content.strip()

            draft_resource = make_draft_resource(
                "article",
                f"实操案例：{user_message[:40]}",
                {"text": content},
            )

            self.emit_step(workflow_outputs, "completed", "生成实操案例", {
                "content": content,
                "sub_steps": ["✅ 案例内容生成完成", "✅ 已生成资源草稿"],
                "draft_resource": draft_resource,
            }, step_id)

            return SkillResult(success=True, data={"content": content, "type": "practice_case", "draft_resource": draft_resource}, summary="实操案例生成完成")
        except Exception as e:
            self.emit_step(workflow_outputs, "completed", "生成实操案例", {
                "content": f"生成失败: {str(e)}",
                "sub_steps": [f"❌ {str(e)}"],
            }, step_id)
            return SkillResult(success=False, error=str(e))


class VideoSearchSkill(BaseSkill):
    """视频检索 — 从B站搜索教学视频"""
    name = "video_search"
    description = "视频检索：从B站搜索相关教学视频，适用于需要视频学习资源的任务"
    icon = "🎬"

    @staticmethod
    def _format_duration(seconds: int) -> str:
        """将秒数格式化为 mm:ss 或 hh:mm:ss"""
        if seconds <= 0:
            return ""
        h = seconds // 3600
        m = (seconds % 3600) // 60
        s = seconds % 60
        if h > 0:
            return f"{h}:{m:02d}:{s:02d}"
        return f"{m}:{s:02d}"

    @staticmethod
    def _format_play(n: int) -> str:
        if n >= 10000:
            return f"{n / 10000:.1f}万"
        return str(n)

    async def execute(self, context: dict, workflow_outputs: list) -> SkillResult:
        from services.bilibili_video_service import search_bilibili_videos

        user_message = context.get("user_message", "")
        ad = context.get("all_modules_data", {})
        search_keywords = ad.get("search_keywords", user_message)

        step_id = self.emit_step(workflow_outputs, "running", "搜索教学视频", {
            "content": f"正在B站搜索: {search_keywords[:30]}...",
            "sub_steps": [],
            "render_type": "video_cards",
        })

        search_result = await search_bilibili_videos(search_keywords, per_keyword=6, total_limit=6)
        videos = search_result.get("videos", [])
        failures = search_result.get("failures", [])
        sub_steps = []
        if videos:
            sub_steps.append(f"✅ 找到 {len(videos)} 个可直达视频")
        if failures:
            sub_steps.append(f"⚠️ {len(failures)} 个关键词未解析到具体视频")
        if not videos:
            sub_steps.append("⚠️ 未找到可直达播放的视频，请换更具体的关键词重试")

        draft_resource = make_draft_resource(
            "video",
            f"视频推荐：{search_keywords[:40]}",
            {"videos": videos, "failures": failures, "query": search_keywords},
        )
        self.emit_step(workflow_outputs, "completed", "搜索教学视频", {
            "content": json.dumps(videos, ensure_ascii=False),
            "sub_steps": sub_steps,
            "render_type": "video_cards",
            "draft_resource": draft_resource if videos else None,
        }, step_id)

        return SkillResult(
            success=True,
            data={"videos": videos, "failures": failures, "type": "video", "draft_resource": draft_resource if videos else None},
            summary=f"找到 {len(videos)} 个B站可直达教学视频",
        )


PPT_PROMPT = """你是一个专业的课件设计专家。根据学生画像和知识点，生成一份结构精美的PPT课件内容。

学生画像：{profile}
主题：{topic}

返回JSON格式：
{{
  "title": "课件标题（精简5-12字）",
  "slides": [
    {{
      "title": "幻灯片标题",
      "content": ["要点1：具体可讲解的内容", "要点2"],
      "notes": "讲师备注：讲解要点或补充信息"
    }}
  ]
}}

设计要求：
- 第1张作为封面（课件标题+副标题概括），最后1张做总结
- 总页数 6-10 张，每页 2-4 条内容要点
- 要点做到简明有深度，适合大学生课堂
- 逻辑递进：引入 → 概念/原理 → 方法/案例 → 应用 → 总结

只返回JSON，不要其他内容。"""


class PptGenSkill(BaseSkill):
    """PPT课件生成 - 支持 .pptx 文件导出"""
    name = "ppt_gen"
    description = "PPT课件生成：根据知识点生成结构化课件，支持导出.pptx文件，适用于需要教学演示的任务"
    icon = "📊"

    async def execute(self, context: dict, workflow_outputs: list) -> SkillResult:
        import os as _os
        from core.database import SessionLocal
        from models.resource import LearningResource
        from services.rag_service import index_resource
        from services.config_service import is_configured as _is_configured

        user_message = context.get("user_message", "")
        user_id = context.get("user_id", "")
        ad = context.get("all_modules_data", {}) or {}

        def normalize_ppt_topic(value: str) -> str:
            text = str(value or "").strip().strip("\"'“”‘’")
            for token in ["PPT", "ppt", "课件", "生成", "制作", "帮我", "请帮我", "给我", "一个", "一份", "助我理解"]:
                text = text.replace(token, "")
            for token in ["的重点是什么", "重点是什么", "重点"]:
                text = text.replace(token, "")
            text = " ".join(text.replace("：", " ").replace(":", " ").split())
            parts = [part for part in text.split(" ") if part]
            if len(parts) == 2:
                text = f"{parts[0]}与{parts[1]}"
            elif len(parts) > 2:
                text = f"{parts[0]}与{''.join(parts[1:])}"
            return text[:60].strip(" ，。；、")

        topic_candidates = [
            ad.get("search_keywords"),
            ad.get("topic"),
            ad.get("code_desc"),
            user_message,
        ]
        ppt_topic = ""
        for candidate in topic_candidates:
            normalized = normalize_ppt_topic(str(candidate or ""))
            if normalized and normalized.lower() not in {"ppt", "ppt课件"}:
                ppt_topic = normalized
                break
        ppt_topic = ppt_topic or user_message

        course_name = str(context.get("course_name") or "").strip()
        knowledge_points = [str(kp).strip() for kp in (context.get("knowledge_points") or []) if str(kp).strip()]
        resolved_binding = None
        if not course_name or not knowledge_points:
            try:
                from services.ppt_model_service import resolve_ppt_binding
                resolved_binding = resolve_ppt_binding(user_id or "default", ppt_topic)
                if resolved_binding.get("status") == "auto_bound":
                    binding = resolved_binding.get("binding") or {}
                    course_name = str(binding.get("course_name") or "").strip()
                    knowledge_points = [
                        str(kp).strip()
                        for kp in (binding.get("knowledge_points") or [])
                        if str(kp).strip()
                    ]
            except Exception:
                resolved_binding = None
        step_id = self.emit_step(workflow_outputs, "running", "创建 AiPPT 分步会话", {
            "sub_steps": ["⏳ 正在创建 AiPPT 工作台会话..."],
        })
        if not course_name or not knowledge_points:
            message = "PPT 需要进入 AiPPT 分步工作台继续配置。你可以选择课程/知识点绑定，或确认创建不参与掌握度更新的拓展课件。"
            pending_session = {
                "pending_binding": True,
                "topic": ppt_topic,
                "original_task": user_message,
                "resolve_status": (resolved_binding or {}).get("status") or "extension_confirm",
                "candidates": (resolved_binding or {}).get("candidates") or [],
                "scope": "pending",
            }
            self.emit_step(workflow_outputs, "completed", "创建 AiPPT 分步会话", {
                "content": json.dumps({
                    "title": ppt_topic or "PPT课件",
                    "ppt_session": pending_session,
                    "message": message,
                }, ensure_ascii=False),
                "sub_steps": [f"✓ {message}"],
                "render_type": "ppt_session",
            }, step_id)
            return SkillResult(
                success=True,
                data={"ppt_session": pending_session, "type": "ppt", "status": "pending_binding"},
                summary="待进入 AiPPT 分步工作台配置",
            )
        try:
            from services.ppt_model_service import create_ppt_session
            ppt_session = await create_ppt_session(
                user_id=user_id or "default",
                topic=ppt_topic,
                course_name=course_name,
                knowledge_points=knowledge_points,
            )
            self.emit_step(workflow_outputs, "completed", "创建 AiPPT 分步会话", {
                "content": json.dumps({
                    "title": ppt_topic or "PPT课件",
                    "ppt_session": ppt_session,
                    "message": "请在 AiPPT 工作台确认大纲和模板后生成 PPT。",
                }, ensure_ascii=False),
                "sub_steps": ["✅ 已创建分步生成会话", "✅ 等待用户确认大纲和模板"],
                "render_type": "ppt_session",
            }, step_id)
            return SkillResult(
                success=True,
                data={"ppt_session": ppt_session, "type": "ppt"},
                summary="已创建 AiPPT 分步生成会话",
            )
        except Exception as exc:
            self.emit_step(workflow_outputs, "completed", "创建 AiPPT 分步会话", {
                "content": f"创建失败: {exc}",
                "sub_steps": [f"❌ 创建 AiPPT 会话失败: {exc}"],
            }, step_id)
            return SkillResult(success=False, error=str(exc))

        step_id = self.emit_step(workflow_outputs, "running", "生成PPT课件", {
            "sub_steps": ["⏳ 正在调用模型生成课件内容..."],
        })

        # 构建学生画像文本
        profile_text = "暂无学生画像"
        profile = context.get("profile")
        if profile:
            import json as _json
            parts = [
                f"专业：{getattr(profile, 'major', '未知')}",
                f"年级：{getattr(profile, 'grade', '未知')}",
            ]
            kb = getattr(profile, 'knowledge_base', {})
            if kb:
                parts.append(f"知识基础：{_json.dumps(kb, ensure_ascii=False)}")
            goal = getattr(profile, 'learning_goal', '')
            if goal:
                parts.append(f"学习目标：{goal}")
            profile_text = "；".join(parts)

        # 1. 调用 LLM 生成JSON内容
        try:
            from core.llm_client import ppt_completion, chat_completion
            use_ppt_model = _is_configured("ppt")
            completion_fn = ppt_completion if use_ppt_model else chat_completion

            resp = await completion_fn([
                {"role": "user", "content": PPT_PROMPT.format(profile=profile_text, topic=user_message)},
            ], temperature=0.5)

            raw = resp.choices[0].message.content.strip()
            if raw.startswith("```"):
                raw = raw.strip("`").strip()
                if raw.startswith("json"):
                    raw = raw[4:].strip()

            ppt_data = json.loads(raw)
            if not isinstance(ppt_data, dict) or "slides" not in ppt_data:
                ppt_data = {"title": user_message, "slides": []}

            sub_steps = [f"✅ LLM生成完成，共 {len(ppt_data.get('slides', []))} 页"]
            self.emit_step(workflow_outputs, "running", "生成PPT课件", {"sub_steps": sub_steps + ["⏳ 正在生成 .pptx 文件..."]}, step_id)
        except Exception as e:
            self.emit_step(workflow_outputs, "completed", "生成PPT课件", {
                "content": f"生成失败: {str(e)}",
                "sub_steps": [f"❌ LLM调用失败: {str(e)}"],
            }, step_id)
            return SkillResult(success=False, error=str(e))

        # 2. 生成 .pptx 文件
        pptx_filename = ""
        pptx_path = ""
        try:
            pptx_filename = self._generate_pptx_file(ppt_data)
            pptx_path = _os.path.join(_os.path.dirname(_os.path.dirname(__file__)), "static", "ppt", pptx_filename)
            sub_steps.append(f"✅ .pptx 文件已生成")
        except Exception as e:
            sub_steps.append(f"⚠️ .pptx 生成失败: {str(e)}")
        self.emit_step(
            workflow_outputs,
            "running",
            "生成PPT课件",
            {"sub_steps": sub_steps + (["⏳ 正在保存到资源库..."] if persist else ["⏳ 正在准备资源草稿..."])},
            step_id,
        )

        # 3. 保存到数据库
        db_id = None
        draft_resource = None
        download_url = f"/static/ppt/{pptx_filename}" if pptx_filename else ""
        if persist:
            try:
                db = SessionLocal()
                try:
                    resource = LearningResource(
                        user_id=user_id,
                        resource_type="ppt",
                        title=ppt_data.get("title", user_message),
                        content={"slides": ppt_data.get("slides", []), "title": ppt_data.get("title", ""), "pptx_file": pptx_filename},
                        tags=["ppt"],
                    )
                    db.add(resource)
                    db.flush()
                    db.commit()
                    db_id = resource.id
                    index_resource(resource.id, user_id or "", json.dumps(ppt_data, ensure_ascii=False)[:4000], "ppt")
                    sub_steps.append(f"✅ 已保存至学习资源库")
                finally:
                    db.close()
            except Exception as e:
                sub_steps.append(f"⚠️ 数据库保存失败: {str(e)}")
        else:
            draft_content = {
                "slides": ppt_data.get("slides", []),
                "title": ppt_data.get("title", ""),
                "pptx_file": pptx_filename,
                "pptx_url": download_url,
            }
            draft_resource = make_draft_resource("ppt", ppt_data.get("title", user_message), draft_content)
            sub_steps.append("✅ 已生成资源草稿")

        # 4. 构建前端渲染数据

        render_content = json.dumps({
            "title": ppt_data.get("title", user_message),
            "slides": ppt_data.get("slides", []),
            "pptx_url": download_url,
            "slide_count": len(ppt_data.get("slides", [])),
            "db_id": db_id,
        }, ensure_ascii=False)

        self.emit_step(workflow_outputs, "completed", "生成PPT课件", {
            "content": render_content,
            "sub_steps": sub_steps,
            "render_type": "ppt_viewer",
            "draft_resource": draft_resource,
        }, step_id)

        return SkillResult(
            success=True,
            data={
                "ppt_json": ppt_data,
                "pptx_filename": pptx_filename,
                "pptx_url": download_url,
                "db_id": db_id,
                "type": "ppt",
                "draft_resource": draft_resource,
            },
            summary=f"PPT课件生成完成 ({len(ppt_data.get('slides', []))} 页)" + (f", .pptx已导出" if pptx_filename else ""),
        )

    def _generate_pptx_file(self, ppt_data: dict) -> str:
        """将 JSON 幻灯片数据转换为 .pptx 文件，返回文件名"""
        import uuid as _uuid
        import os as _os
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        slides_data = ppt_data.get("slides", [])
        title_text = ppt_data.get("title", "课件")

        # 封面
        if slides_data:
            self._make_cover_slide(prs, title_text, slides_data[0])
            content_slides = slides_data[1:]
        else:
            content_slides = []

        # 内容页
        for sd in content_slides:
            self._make_content_slide(prs, sd.get("title", ""), sd.get("content", []), sd.get("notes", ""))

        # 如果没有内容，至少加一张空页
        if len(prs.slides) == 0:
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = title_text

        static_dir = _os.path.join(_os.path.dirname(_os.path.dirname(__file__)), "static", "ppt")
        _os.makedirs(static_dir, exist_ok=True)

        filename = f"{_uuid.uuid4().hex[:8]}.pptx"
        filepath = _os.path.join(static_dir, filename)
        prs.save(filepath)
        return filename

    def _make_cover_slide(self, prs, title: str, first_slide: dict):
        """制作封面页"""
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        slide_layout = prs.slide_layouts[6]  # blank
        slide = prs.slides.add_slide(slide_layout)

        # 背景色块
        left, top = Inches(0), Inches(0)
        width, height = prs.slide_width, prs.slide_height
        bg_shape = slide.shapes.add_shape(
            1, left, top, width, height  # MSO_SHAPE.RECTANGLE
        )
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = RGBColor(0x1A, 0x36, 0x5D)
        bg_shape.line.fill.background()

        # 标题
        txBox = slide.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(10.3), Inches(1.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.alignment = PP_ALIGN.CENTER

        # 副标题
        subtitle = first_slide.get("title", "")
        if subtitle:
            txBox2 = slide.shapes.add_textbox(Inches(1.5), Inches(3.8), Inches(10.3), Inches(1.0))
            tf2 = txBox2.text_frame
            p2 = tf2.paragraphs[0]
            p2.text = subtitle
            p2.font.size = Pt(24)
            p2.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
            p2.alignment = PP_ALIGN.CENTER

        # 装饰线
        line = slide.shapes.add_shape(
            1, Inches(4.5), Inches(5.2), Inches(2.8), Inches(0.04)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(0x40, 0x9E, 0xFF)
        line.line.fill.background()

    def _make_content_slide(self, prs, title: str, content: list, notes: str):
        """制作内容页"""
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        slide_layout = prs.slide_layouts[6]  # blank
        slide = prs.slides.add_slide(slide_layout)

        # 顶部装饰条
        top_bar = slide.shapes.add_shape(
            1, Inches(0), Inches(0), prs.slide_width, Inches(0.08)
        )
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = RGBColor(0x40, 0x9E, 0xFF)
        top_bar.line.fill.background()

        # 标题
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.9))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x1A, 0x36, 0x5D)

        # 内容要点
        content_top = Inches(1.8)
        txBox2 = slide.shapes.add_textbox(Inches(1.2), content_top, Inches(10.9), Inches(5.0))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True

        for i, point in enumerate(content):
            if i == 0:
                p2 = tf2.paragraphs[0]
            else:
                p2 = tf2.add_paragraph()
            p2.text = f"• {point}"
            p2.font.size = Pt(20)
            p2.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
            p2.space_after = Pt(14)
            p2.space_before = Pt(4)

        # 讲师备注
        if notes:
            txBox3 = slide.shapes.add_textbox(Inches(1.2), Inches(6.2), Inches(10.9), Inches(0.9))
            tf3 = txBox3.text_frame
            tf3.word_wrap = True
            p3 = tf3.paragraphs[0]
            p3.text = f"📝 {notes}"
            p3.font.size = Pt(13)
            p3.font.italic = True
            p3.font.color.rgb = RGBColor(0x90, 0x93, 0x99)

        # 页码
        slide_num = len(prs.slides)
        txBox4 = slide.shapes.add_textbox(Inches(12.0), Inches(6.9), Inches(1.0), Inches(0.4))
        tf4 = txBox4.text_frame
        p4 = tf4.paragraphs[0]
        p4.text = str(slide_num)
        p4.font.size = Pt(11)
        p4.font.color.rgb = RGBColor(0xC0, 0xC4, 0xCC)
        p4.alignment = PP_ALIGN.RIGHT


# ============================================================
# 初始化：注册所有内置 skills
# ============================================================

def init_skills():
    """注册所有内置 skills，在应用启动时调用"""
    builtin_skills = [
        DeepSearchSkill(),
        CodeAnalysisSkill(),
        MindmapSkill(),
        QuizGenSkill(),
        VideoSearchSkill(),
        PptGenSkill(),
        ArticleGenSkill(),
        CodeGenSkill(),
        PracticeCaseSkill(),
    ]
    for skill in builtin_skills:
        if skill.name not in _SKILL_REGISTRY:
            register_skill(skill)
